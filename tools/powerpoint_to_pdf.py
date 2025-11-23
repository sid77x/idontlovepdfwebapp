"""Convert PowerPoint (PPTX) to PDF."""
import os
import sys
import streamlit as st
from pptx import Presentation
from reportlab.lib.pagesizes import landscape, letter
from reportlab.pdfgen import canvas
from PIL import Image
from io import BytesIO
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from utils.file_ops import save_uploaded_file, get_unique_filename, get_file_size_mb
from utils.preview import show_pdf_preview


def powerpoint_to_pdf(input_path: str, output_path: str) -> bool:
    """
    Convert PowerPoint to PDF (simplified version - text only).
    
    Args:
        input_path: Input PPTX path
        output_path: Output PDF path
    
    Returns:
        bool: Success status
    """
    try:
        # Open PowerPoint
        prs = Presentation(input_path)
        
        # Create PDF with landscape orientation
        page_width, page_height = landscape(letter)
        c = canvas.Canvas(output_path, pagesize=landscape(letter))
        
        # Process each slide
        for slide_num, slide in enumerate(prs.slides):
            if slide_num > 0:
                c.showPage()  # New page for each slide
            
            # Add slide number
            c.setFont("Helvetica", 10)
            c.drawString(30, 30, f"Slide {slide_num + 1}")
            
            # Extract and position text
            y_position = page_height - 100
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Determine font size based on shape type
                    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or hasattr(shape, 'text_frame'):
                        # Check if it's likely a title (larger font)
                        if y_position > page_height - 150:
                            font_size = 24
                            c.setFont("Helvetica-Bold", font_size)
                        else:
                            font_size = 14
                            c.setFont("Helvetica", font_size)
                        
                        # Split text into lines to fit page width
                        max_width = page_width - 100
                        words = text.split()
                        lines = []
                        current_line = []
                        
                        for word in words:
                            current_line.append(word)
                            line_text = ' '.join(current_line)
                            if c.stringWidth(line_text, "Helvetica", font_size) > max_width:
                                if len(current_line) > 1:
                                    current_line.pop()
                                    lines.append(' '.join(current_line))
                                    current_line = [word]
                                else:
                                    lines.append(line_text)
                                    current_line = []
                        
                        if current_line:
                            lines.append(' '.join(current_line))
                        
                        # Draw lines
                        for line in lines:
                            if y_position < 100:
                                break
                            c.drawString(50, y_position, line)
                            y_position -= font_size + 5
                        
                        y_position -= 10  # Extra space after each text block
        
        # Save PDF
        c.save()
        
        return True
        
    except Exception as e:
        st.error(f"Error converting PowerPoint to PDF: {str(e)}")
        return False


def render_powerpoint_to_pdf_tool():
    """Render the PowerPoint to PDF tool interface."""
    st.title("📊 PowerPoint to PDF")
    st.write("Convert your PowerPoint presentations to PDF format")
    
    # File upload
    uploaded_file = st.file_uploader(
        "Choose a PowerPoint file",
        type=['pptx'],
        help="Select the PowerPoint file you want to convert to PDF"
    )
    
    if uploaded_file:
        # Display file info
        file_size = get_file_size_mb(uploaded_file)
        st.info(f"📊 **{uploaded_file.name}** ({file_size:.2f} MB)")
        
        # Save uploaded file
        input_path = save_uploaded_file(uploaded_file, "powerpoint_to_pdf")
        
        if input_path:
            try:
                # Load presentation to get info
                prs = Presentation(input_path)
                num_slides = len(prs.slides)
                
                st.success(f"✅ PowerPoint loaded: {num_slides} slides")
                
                st.warning("⚠️ **Note:** This is a simplified converter. Images, charts, and complex formatting may not be preserved.")
                st.info("💡 **Tip:** For best results with complex presentations, use PowerPoint's built-in 'Save as PDF' feature.")
                
                # Convert button
                if st.button("🔄 Convert to PDF", type="primary", use_container_width=True):
                    with st.spinner("Converting PowerPoint to PDF..."):
                        # Create output path
                        output_dir = os.path.join("output", "powerpoint_to_pdf")
                        os.makedirs(output_dir, exist_ok=True)
                        output_path = os.path.join(output_dir, get_unique_filename("converted.pdf"))
                        
                        # Convert PowerPoint to PDF
                        success = powerpoint_to_pdf(input_path, output_path)
                        
                        if success:
                            st.success(f"✅ Successfully converted PowerPoint to PDF!")
                            
                            # Show preview
                            st.write("### 🔍 Preview")
                            try:
                                with open(output_path, "rb") as f:
                                    preview_data = f.read()
                                
                                show_pdf_preview(preview_data, height=1000)
                            except Exception as e:
                                st.warning(f"Preview unavailable: {str(e)}")
                            
                            # Show file info
                            output_size = os.path.getsize(output_path) / (1024 * 1024)
                            st.info(f"📄 Output PDF size: {output_size:.2f} MB")
                            
                            # Download button
                            with open(output_path, "rb") as f:
                                st.download_button(
                                    label="⬇️ Download PDF",
                                    data=f,
                                    file_name="converted.pdf",
                                    mime="application/pdf",
                                    use_container_width=True
                                )
                
            except Exception as e:
                st.error(f"❌ Error processing PowerPoint file: {str(e)}")
    
    else:
        # Help section
        st.write("### 📖 How to use")
        st.markdown("""
        1. **Upload** your PowerPoint file (.pptx)
        2. **Click** Convert to PDF
        3. **Preview** the result
        4. **Download** your PDF file
        
        **What to expect:**
        - ✅ Text content extracted
        - ✅ Slide-by-slide conversion
        - ✅ Landscape orientation
        - ⚠️ Simplified formatting
        
        **Supported:**
        - Microsoft PowerPoint 2007+ (.pptx)
        - Google Slides (download as .pptx first)
        - LibreOffice Impress (.pptx)
        
        **Limitations:**
        - ⚠️ Images not included
        - ⚠️ Charts not included
        - ⚠️ Animations ignored
        - ⚠️ Complex layouts simplified
        
        **Best for:**
        - Text-heavy presentations
        - Quick conversions
        - Extracting text content
        
        **For complex presentations:**
        Use PowerPoint's native "File → Save As → PDF" for best results.
        """)


if __name__ == "__main__":
    render_powerpoint_to_pdf_tool()
