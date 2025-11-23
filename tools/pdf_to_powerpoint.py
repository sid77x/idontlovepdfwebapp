"""Convert PDF to PowerPoint (PPTX)."""
import os
import sys
import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from io import BytesIO

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from utils.file_ops import save_uploaded_file, get_unique_filename, get_file_size_mb


def pdf_to_powerpoint(input_path: str, output_path: str, dpi: int = 150) -> bool:
    """
    Convert PDF to PowerPoint by converting each page to an image.
    
    Args:
        input_path: Input PDF path
        output_path: Output PPTX path
        dpi: Resolution for page rendering
    
    Returns:
        bool: Success status
    """
    try:
        # Open PDF
        pdf_document = fitz.open(input_path)
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Set slide dimensions to match typical presentation (16:9)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Process each page
        zoom = dpi / 72
        matrix = fitz.Matrix(zoom, zoom)
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            # Render page to image
            pix = page.get_pixmap(matrix=matrix)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            # Save image to BytesIO
            img_buffer = BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Calculate image dimensions to fit slide
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            img_width = img.width
            img_height = img.height
            
            # Calculate scaling to fit slide while maintaining aspect ratio
            width_ratio = slide_width / img_width
            height_ratio = slide_height / img_height
            scale_ratio = min(width_ratio, height_ratio)
            
            new_width = int(img_width * scale_ratio)
            new_height = int(img_height * scale_ratio)
            
            # Center the image
            left = (slide_width - new_width) / 2
            top = (slide_height - new_height) / 2
            
            # Add image to slide
            slide.shapes.add_picture(img_buffer, left, top, width=new_width, height=new_height)
        
        # Save presentation
        prs.save(output_path)
        pdf_document.close()
        
        return True
        
    except Exception as e:
        st.error(f"Error converting PDF to PowerPoint: {str(e)}")
        return False


def render_pdf_to_powerpoint_tool():
    """Render the PDF to PowerPoint tool interface."""
    st.title("📊 PDF to PowerPoint")
    st.write("Convert your PDF pages to PowerPoint slides")
    
    # File upload
    uploaded_file = st.file_uploader(
        "Choose a PDF file",
        type=['pdf'],
        help="Select the PDF file you want to convert to PowerPoint"
    )
    
    if uploaded_file:
        # Display file info
        file_size = get_file_size_mb(uploaded_file)
        st.info(f"📄 **{uploaded_file.name}** ({file_size:.2f} MB)")
        
        # Save uploaded file
        input_path = save_uploaded_file(uploaded_file, "pdf_to_powerpoint")
        
        if input_path:
            try:
                # Get page count
                pdf_doc = fitz.open(input_path)
                total_pages = len(pdf_doc)
                pdf_doc.close()
                
                st.success(f"✅ PDF loaded: {total_pages} pages")
                
                # Conversion settings
                st.write("### ⚙️ Conversion Settings")
                
                dpi = st.slider(
                    "Image Quality (DPI)",
                    min_value=72,
                    max_value=200,
                    value=150,
                    step=10,
                    help="Higher DPI = better quality but larger file size"
                )
                
                st.info("💡 **Note:** Each PDF page will become one PowerPoint slide as an image.")
                st.warning("⚠️ Text will not be editable in PowerPoint - pages are converted to images.")
                
                # Convert button
                if st.button("🔄 Convert to PowerPoint", type="primary", use_container_width=True):
                    with st.spinner("Converting PDF to PowerPoint... This may take a moment."):
                        # Create output path
                        output_dir = os.path.join("output", "pdf_to_powerpoint")
                        os.makedirs(output_dir, exist_ok=True)
                        output_path = os.path.join(output_dir, get_unique_filename("converted.pptx"))
                        
                        # Convert PDF to PowerPoint
                        success = pdf_to_powerpoint(input_path, output_path, dpi)
                        
                        if success:
                            st.success(f"✅ Successfully converted PDF to PowerPoint!")
                            
                            # Show file info
                            output_size = os.path.getsize(output_path) / (1024 * 1024)
                            st.info(f"📊 Output file: {output_size:.2f} MB • {total_pages} slides")
                            
                            # Download button
                            with open(output_path, "rb") as f:
                                st.download_button(
                                    label="⬇️ Download PowerPoint",
                                    data=f,
                                    file_name="converted.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    use_container_width=True
                                )
                            
                            st.success("💡 **Tip:** Open in PowerPoint, Google Slides, or LibreOffice Impress.")
                
            except Exception as e:
                st.error(f"❌ Error processing PDF: {str(e)}")
    
    else:
        # Help section
        st.write("### 📖 How to use")
        st.markdown("""
        1. **Upload** your PDF file
        2. **Adjust** image quality if needed
        3. **Click** Convert to PowerPoint
        4. **Download** your presentation
        
        **What to expect:**
        - Each PDF page = One slide
        - Pages converted to high-quality images
        - Slides sized for standard 16:9 presentations
        - Images centered on slides
        
        **Best for:**
        - Converting presentations back to editable format
        - Sharing PDF slides as PowerPoint
        - Creating slide decks from documents
        
        **Limitations:**
        - ⚠️ Text is not editable (image-based)
        - ⚠️ No animations or transitions
        - ⚠️ Larger file sizes than original PDF
        
        **Recommended DPI:**
        - 72-100: Screen viewing only
        - 150: Good balance (recommended)
        - 200: High quality for printing
        """)


if __name__ == "__main__":
    render_pdf_to_powerpoint_tool()
